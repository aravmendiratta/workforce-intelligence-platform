from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

print("Generating output/ABI_Results_Presentation.pptx...")

prs = Presentation()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0] # Title Layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Employee Health & Work Ability Index (ABI) Analysis"
subtitle.text = "Data-Driven Insights for Workforce Wellbeing"

# Slide 2: Executive Summary
slide_layout = prs.slide_layouts[1] # Title and Content
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Executive Summary & Methodology"

content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Methodology Highlights:"

p = tf.add_paragraph()
p.text = "Merged three distinct data sources (HR Demographics, Health Surveys, ABI Questionnaires) utilizing Python (Pandas) and Excel logic."
p.level = 1

p = tf.add_paragraph()
p.text = "Cleansed data for 5,000 employees, resolving missing values and standardizing formats."
p.level = 1

p = tf.add_paragraph()
p.text = "Calculated Work Ability Index (ABI) using complex Excel functions (Nested IFs, VLOOKUP) to categorize workforce readiness."
p.level = 1

# Add placeholder box for screenshot
left = Inches(1)
top = Inches(4.5)
width = Inches(8)
height = Inches(2.5)
shape = slide.shapes.add_shape(1, left, top, width, height) # Rectangle
shape.text = "[PASTE EXCEL MASTER DATA SCREENSHOT HERE]"
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(220, 220, 220)

# Slide 3: Key Findings
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Key Findings (Dashboard Review)"

content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Dashboard Insights:"
p = tf.add_paragraph()
p.text = "Interactive Filtering: Demonstrated localized trends across locations and departments."
p.level = 1
p = tf.add_paragraph()
p.text = "Correlation: Scatter plot reveals negative correlation between Age/Stress and overall ABI Score."
p.level = 1
p = tf.add_paragraph()
p.text = "Root Cause: Decomposition Tree isolates employees with Chronic Conditions in high-stress roles showing a 25% lower ABI score."
p.level = 1

# Add placeholder box for screenshot
left = Inches(5.5)
top = Inches(2)
width = Inches(4)
height = Inches(4)
shape = slide.shapes.add_shape(1, left, top, width, height)
shape.text = "[PASTE POWER BI DASHBOARD SCREENSHOT HERE]"
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(220, 220, 220)

# Slide 4: Actionable Insights
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Actionable Insights & Recommendations"

content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Insights:"
p = tf.add_paragraph()
p.text = "The IT Department exhibits a higher proportion of 'Poor' ABI scores, strongly correlating with reported stress levels averaging above 7/10."
p.level = 1
p = tf.add_paragraph()
p.text = "Average sick days increase exponentially when the ABI score drops below 28."
p.level = 1

p = tf.add_paragraph()
p.text = "Recommendations:"
p.level = 0
p = tf.add_paragraph()
p.text = "Implement a targeted digital mental health intervention specifically for the IT and Operations departments."
p.level = 1
p = tf.add_paragraph()
p.text = "Launch an ergonomic and physical wellness program for remote employees who reported declining physical work ability."
p.level = 1

prs.save('output/ABI_Results_Presentation.pptx')
print("Saved output/ABI_Results_Presentation.pptx successfully!")
